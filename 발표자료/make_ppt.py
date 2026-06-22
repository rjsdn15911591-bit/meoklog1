"""
먹로그 발표 PPT 생성 스크립트 v2
AI와 머신러닝 기말고사 과제 발표 자료 — 10분 분량, 16:9
폰트: 본문=NanumMyeongjo / 제목·UI=KERIS KEDU B(bold) · KERIS KEDU R(regular)
"""

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── 색상 팔레트 (먹로그 디자인 시스템) ──────────────────────────────
C_CANVAS   = RGBColor(0xF9, 0xF6, 0xF1)
C_COBALT   = RGBColor(0x50, 0x58, 0xF0)
C_SAGE     = RGBColor(0x70, 0xB0, 0x80)
C_PEACH    = RGBColor(0xFF, 0xB0, 0x84)
C_OCHRE    = RGBColor(0xE8, 0xB9, 0x4A)
C_CORAL    = RGBColor(0xE8, 0x5D, 0x4A)
C_LAVEN    = RGBColor(0xC0, 0xC0, 0xF0)
C_TEAL     = RGBColor(0x90, 0xC0, 0xC0)
C_INK      = RGBColor(0x1A, 0x1A, 0x1A)
C_BODY     = RGBColor(0x3A, 0x3A, 0x3A)
C_MUTED    = RGBColor(0x6A, 0x6A, 0x6A)
C_MUTED_S  = RGBColor(0x9A, 0x9A, 0x9A)
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_HAIRLINE = RGBColor(0xE0, 0xD8, 0xCC)
C_SURF     = RGBColor(0xF2, 0xED, 0xE4)

# ── 폰트 정의 ─────────────────────────────────────────────────────────
# 케리스 케듀체: 감정/UI 레이어 (제목, 챕터명, 앱명, 카드 제목)
FONT_TITLE = "KERIS KEDU B"    # Bold — 슬라이드 제목, 카드 제목, "먹로그"
FONT_UI    = "KERIS KEDU R"    # Regular — 챕터명, 부제목, 레이블
# 나눔명조: 데이터/정보 레이어 (본문, 수치, 기술 내용)
FONT_BODY  = "NanumMyeongjo"   # 본문 전용

# ── 슬라이드 치수 (16:9) ─────────────────────────────────────────────
SW = Cm(33.867)
SH = Cm(19.05)

# ── 레이아웃 상수 ────────────────────────────────────────────────────
HDR_Y = Cm(0);     HDR_H = Cm(1.45)
TTL_Y = Cm(1.6);   TTL_H = Cm(1.55)
DIV_Y = Cm(3.25);  DIV_H = Cm(0.07)
CNT_Y = Cm(3.55);  CNT_H = Cm(14.4)
FTR_Y = Cm(18.05); FTR_H = Cm(1.0)
ML    = Cm(1.8)
MR    = Cm(1.8)
CW    = SW - ML - MR

LOGO_PATH    = r"c:\Desktop\문서\AI_GUNWOO\Muklog\app logo\KakaoTalk_20260620_204355294_03.png"
TOTAL_SLIDES = 13

# ═══════════════════════════════════════════════════════════════════════
# 유틸 함수
# ═══════════════════════════════════════════════════════════════════════

def prs_new():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color=C_CANVAS):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, x, y, w, h, fill_color, line_color=None, line_w=None):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color:
        shp.line.color.rgb = line_color
        if line_w: shp.line.width = Pt(line_w)
    else:
        shp.line.fill.background()
    return shp

def tbx(slide, x, y, w, h):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    return box, tf

def para(tf, text, size=12, bold=False, color=C_BODY,
         align=PP_ALIGN.LEFT, space_before=0, space_after=0,
         italic=False, font=None):
    """tf에 단락 추가. font 미지정 시 FONT_BODY 사용."""
    if font is None:
        font = FONT_BODY
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == '':
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    run = p.add_run()
    run.text = text
    run.font.name  = font
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    p.alignment = align
    if space_before: p.space_before = Pt(space_before)
    if space_after:  p.space_after  = Pt(space_after)
    return p

def para_multi(tf, lines):
    """
    lines: list of dict
      {'text','size','bold','color','align','space_before','space_after','italic','font'}
    font 미지정 시 FONT_BODY 사용.
    """
    first = True
    for ln in lines:
        if isinstance(ln, str):
            ln = {'text': ln}
        t   = ln.get('text', '')
        sz  = ln.get('size', 12)
        bd  = ln.get('bold', False)
        col = ln.get('color', C_BODY)
        al  = ln.get('align', PP_ALIGN.LEFT)
        sb  = ln.get('space_before', 0)
        sa  = ln.get('space_after', 0)
        it  = ln.get('italic', False)
        ft  = ln.get('font', FONT_BODY)
        if first and len(tf.paragraphs) == 1 and tf.paragraphs[0].text == '':
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = t
        run.font.name   = ft
        run.font.size   = Pt(sz)
        run.font.bold   = bd
        run.font.italic = it
        run.font.color.rgb = col
        p.alignment = al
        if sb: p.space_before = Pt(sb)
        if sa: p.space_after  = Pt(sa)

def header(slide, chapter_text, color=C_COBALT):
    """상단 헤더 바 + 챕터명 (KERIS KEDU R)"""
    rect(slide, Cm(0), HDR_Y, SW, HDR_H, color)
    bx, tf = tbx(slide, ML, Cm(0.1), CW, HDR_H)
    para(tf, chapter_text, size=11, bold=False, color=C_WHITE, font=FONT_UI)

def title_block(slide, title_text, subtitle_text=None):
    """제목(KERIS KEDU B) + 부제목(KERIS KEDU R) — 고정 위치"""
    bx, tf = tbx(slide, ML, TTL_Y, CW, TTL_H)
    para(tf, title_text, size=25, bold=False, color=C_INK, font=FONT_TITLE)
    if subtitle_text:
        para(tf, subtitle_text, size=12, color=C_MUTED, space_before=2, font=FONT_UI)

def divider(slide, color=C_PEACH):
    rect(slide, ML, DIV_Y, CW, DIV_H, color)

def footer_bar(slide, num):
    rect(slide, Cm(0), FTR_Y, SW, FTR_H, C_SURF)
    bx, tf = tbx(slide, ML, FTR_Y + Cm(0.15), Cm(10), Cm(0.7))
    para(tf, "팀 먹로그  |  AI와 머신러닝 기말 과제  |  25학번 이건우 · 23학번 윤서영",
         size=8, color=C_MUTED, font=FONT_UI)
    bx2, tf2 = tbx(slide, SW - Cm(4), FTR_Y + Cm(0.15), Cm(3.5), Cm(0.7))
    para(tf2, f"{num}  /  {TOTAL_SLIDES}", size=9, bold=False, color=C_MUTED,
         align=PP_ALIGN.RIGHT, font=FONT_BODY)

def std_frame(slide, chapter, title, subtitle=None, num=1):
    bg(slide)
    header(slide, chapter)
    title_block(slide, title, subtitle)
    divider(slide)
    footer_bar(slide, num)

def card(slide, x, y, w, h, bg_color, title_text, lines,
         title_color=C_WHITE, text_color=C_WHITE, t_size=12, c_size=10.5):
    """색상 카드 — 제목 KERIS KEDU B, 본문 NanumMyeongjo"""
    rect(slide, x, y, w, h, bg_color)
    bx, tf = tbx(slide, x + Cm(0.35), y + Cm(0.2), w - Cm(0.5), Cm(0.65))
    para(tf, title_text, size=t_size, bold=False, color=title_color, font=FONT_TITLE)
    if lines:
        bx2, tf2 = tbx(slide, x + Cm(0.35), y + Cm(0.9), w - Cm(0.5), h - Cm(1.1))
        for ln in lines:
            if isinstance(ln, str):
                para(tf2, ln, size=c_size, color=text_color, font=FONT_BODY)
            else:
                if tf2.paragraphs[0].text != '':
                    p = tf2.add_paragraph()
                else:
                    p = tf2.paragraphs[0]
                run = p.add_run()
                run.text = ln.get('text', '')
                run.font.name  = ln.get('font', FONT_BODY)
                run.font.size  = Pt(ln.get('size', c_size))
                run.font.bold  = ln.get('bold', False)
                run.font.color.rgb = ln.get('color', text_color)

# ═══════════════════════════════════════════════════════════════════════
# 슬라이드 생성
# ═══════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    sl = blank_slide(prs)
    bg(sl, C_COBALT)
    rect(sl, Cm(0), Cm(0), SW * 0.62, SH, C_CANVAS)

    # 앱 이름 — KERIS KEDU B
    bx, tf = tbx(sl, Cm(2.2), Cm(3.2), Cm(16), Cm(3.5))
    para(tf, "먹로그", size=60, bold=False, color=C_COBALT, font=FONT_TITLE)

    # 영문명 — KERIS KEDU R italic
    bx2, tf2 = tbx(sl, Cm(2.2), Cm(6.8), Cm(18), Cm(1.2))
    para(tf2, "MealLog", size=22, bold=False, color=C_MUTED, italic=True, font=FONT_UI)

    # 부제목 — KERIS KEDU B
    bx3, tf3 = tbx(sl, Cm(2.2), Cm(8.4), Cm(18), Cm(2.0))
    para(tf3, "AI 기반 그룹형 식단 관리 웹 앱", size=18, bold=False, color=C_INK, font=FONT_TITLE)
    para(tf3, "사진 한 장으로 칼로리·영양소 자동 분석, 친구와 함께 식단 비교",
         size=12, color=C_MUTED, space_before=4, font=FONT_UI)

    rect(sl, Cm(2.2), Cm(11.2), Cm(14), Cm(0.08), C_PEACH)

    # 과목·팀 정보 — NanumMyeongjo
    bx4, tf4 = tbx(sl, Cm(2.2), Cm(11.6), Cm(16), Cm(3.0))
    para_multi(tf4, [
        {'text': '과목  :  AI와 머신러닝',                       'size': 11, 'color': C_BODY, 'font': FONT_BODY},
        {'text': '발표일 :  2026년 6월 22일',                    'size': 11, 'color': C_BODY, 'font': FONT_BODY, 'space_before': 3},
        {'text': '팀명  :  먹로그',                               'size': 11, 'color': C_BODY, 'font': FONT_BODY, 'space_before': 3},
        {'text': '팀원  :  25학번 이건우  ·  23학번 윤서영',      'size': 11, 'color': C_BODY, 'font': FONT_BODY, 'space_before': 3},
    ])

    if os.path.exists(LOGO_PATH):
        sl.shapes.add_picture(LOGO_PATH, Cm(22.5), Cm(4.5), Cm(8.5), Cm(8.5))

    bx5, tf5 = tbx(sl, Cm(22.0), Cm(16.5), Cm(9), Cm(1.2))
    para(tf5, "v1.6.0  ·  Phase 1 + 2 완성", size=10, color=C_WHITE,
         align=PP_ALIGN.CENTER, font=FONT_BODY)
    return sl


def slide_02_toc(prs):
    sl = blank_slide(prs)
    std_frame(sl, "목차", "발표 순서", num=2)

    items = [
        ("01", "프로젝트 개요",     "문제 정의와 먹로그의 핵심 가치",            C_COBALT),
        ("02", "AI 음식 분석 기술", "GPT-4o Vision 3단계 추론 메커니즘",         C_LAVEN),
        ("03", "주요 기능",         "그룹 공유·비교, 개인 건강 관리",             C_SAGE),
        ("04", "프로그램 시연",     "직접 화면을 보며 기능 설명",                 C_PEACH),
        ("05", "기술 아키텍처",     "Next.js · FastAPI · Supabase 구성",         C_OCHRE),
        ("06", "개발 성과 & 향후",  "Phase 완성 내역 + 스마트폰 앱 출시 계획",    C_CORAL),
    ]

    start_y = CNT_Y + Cm(0.2)
    row_h   = Cm(2.2)
    num_w   = Cm(1.6)
    ttl_w   = Cm(8.5)

    for i, (num_t, ttl, sub, col) in enumerate(items):
        y = start_y + i * row_h
        rect(sl, ML, y, num_w, Cm(1.7), col)
        bx_n, tf_n = tbx(sl, ML, y + Cm(0.2), num_w, Cm(1.3))
        para(tf_n, num_t, size=17, bold=False, color=C_WHITE,
             align=PP_ALIGN.CENTER, font=FONT_TITLE)
        bx_t, tf_t = tbx(sl, ML + num_w + Cm(0.4), y, ttl_w, Cm(1.7))
        para(tf_t, ttl, size=14, bold=False, color=C_INK, font=FONT_TITLE)
        para(tf_t, sub, size=10, color=C_MUTED, space_before=2, font=FONT_UI)
        if i < len(items) - 1:
            rect(sl, ML, y + Cm(1.85), CW, Cm(0.04), C_HAIRLINE)
    return sl


def slide_03_problem(prs):
    sl = blank_slide(prs)
    std_frame(sl, "01.  프로젝트 개요", "왜 식단 기록이 어려운가?",
              subtitle="기존 식단 앱의 한계와 먹로그의 탄생 배경", num=3)

    col_w = (CW - Cm(0.6)) / 2
    problems = [
        (C_CORAL,  "❶  직접 입력의 번거로움",
         ["음식명·칼로리를 일일이 검색해야 함", "하루 3끼 × 3~4가지 = 매일 10회 이상 입력"]),
        (C_OCHRE,  "❷  정확도 문제",
         ["조리법·양에 따라 DB 수치와 괴리 발생", "직접 측정 없이는 오차 20~40% 발생"]),
        (C_COBALT, "❸  혼자만의 기록, 동기 부여 한계",
         ["소셜 요소·경쟁 구조 없이 지속 어려움", "작심삼일로 끝나는 경우가 대부분"]),
        (C_MUTED,  "❹  복잡한 UI · 높은 진입 장벽",
         ["기존 앱은 기능이 많아 처음 쓰기 어려움", "온보딩 완료율 평균 30% 미만"]),
    ]

    card_h = Cm(3.3)
    for i, (col, ttl, lines) in enumerate(problems):
        cx = ML + (i % 2) * (col_w + Cm(0.6))
        cy = CNT_Y + Cm(0.1) + (i // 2) * (card_h + Cm(0.3))
        txt = C_WHITE if col in (C_CORAL, C_COBALT) else C_INK
        card(sl, cx, cy, col_w, card_h, col, ttl,
             [{'text': l, 'size': 10, 'color': txt, 'font': FONT_BODY} for l in lines],
             title_color=txt, text_color=txt, t_size=12)

    rect(sl, ML, CNT_Y + Cm(7.2), CW, Cm(0.07), C_PEACH)
    bx, tf = tbx(sl, ML, CNT_Y + Cm(7.5), CW, Cm(1.2))
    para(tf, '"사진 한 장 찍는 것만큼 쉬워야 한다"  —  먹로그의 출발점',
         size=12, bold=False, color=C_COBALT, font=FONT_TITLE)
    return sl


def slide_04_solution(prs):
    sl = blank_slide(prs)
    std_frame(sl, "01.  프로젝트 개요", "먹로그의 핵심 가치",
              subtitle="네 가지 원칙으로 기존 앱의 한계를 극복", num=4)

    values = [
        (C_COBALT, "사진 한 장으로 끝",
         "직접 검색·입력 불필요. 카메라로 찍으면 AI가 메뉴·칼로리·탄단지를 자동 분석",
         "→ 기록 시간 평균 2분 → 15초로 단축"),
        (C_SAGE,   "AI 3단계 정밀 추론",
         "GPT-4o Vision이 크기 기준점 측정 → 밀도 결정 → 칼로리 계산 순으로 추론",
         "→ 단순 인식이 아닌 물리적 계산 기반 추정"),
        (C_PEACH,  "그룹 비교로 동기 부여",
         "친구·가족과 그룹을 만들어 하루 섭취량 실시간 비교, 이모티콘 반응·댓글",
         "→ 소셜 피드 방식으로 꾸준한 기록 유지"),
        (C_OCHRE,  "개인 맞춤 목표",
         "키·체중·활동량·목표 입력 → Mifflin-St Jeor 공식으로 목표 칼로리 자동 계산",
         "→ AI 코치가 맞춤 식단·운동 루틴 생성"),
    ]

    card_h = Cm(3.15)
    col_w  = (CW - Cm(0.6)) / 2
    for i, (col, ttl, desc1, desc2) in enumerate(values):
        cx = ML + (i % 2) * (col_w + Cm(0.6))
        cy = CNT_Y + Cm(0.1) + (i // 2) * (card_h + Cm(0.35))
        txt = C_WHITE if col in (C_COBALT, C_SAGE) else C_INK
        card(sl, cx, cy, col_w, card_h, col, ttl,
             [{'text': desc1, 'size': 10.5, 'color': txt, 'font': FONT_BODY},
              {'text': desc2, 'size': 9.5,  'color': txt, 'bold': True, 'font': FONT_BODY}],
             title_color=txt, t_size=13)
    return sl


def slide_05_ai_overview(prs):
    sl = blank_slide(prs)
    std_frame(sl, "02.  AI 음식 분석 기술", "GPT-4o Vision 3단계 추론",
              subtitle="단순 이미지 인식을 넘어 물리적 계산으로 칼로리를 도출", num=5)

    steps = [
        (C_LAVEN,  "입력",    "음식 사진",
         ["사용자가 카메라로 촬영", "Cloudinary에 자동 업로드", "GPT-4o Vision에 이미지 전달"]),
        (C_COBALT, "STEP 1",  "크기 기준점 찾기",
         ["사진 속 기준 물체 탐지", "(젓가락·그릇·손·테이블)", "실제 g / ml 수량 추정"]),
        (C_SAGE,   "STEP 2",  "밀도 결정",
         ["조리 상태·재료 판단", "내장 밀도표에서 kcal/100g 선택", "(볶음·국물·튀김별 다름)"]),
        (C_PEACH,  "STEP 3",  "칼로리 계산",
         ["calories = g × kcal/100g ÷ 100", "탄수화물·단백질·지방 비율 산출", "JSON 형식으로 결과 반환"]),
        (C_OCHRE,  "출력",    "결과 확인",
         ["음식명·칼로리·탄단지 표시", "½ ~ 2인분 양 조절 슬라이더", "사용자 확인 후 DB 저장"]),
    ]

    step_w = (CW - Cm(0.4) * 4) / 5
    step_h = Cm(9.5)
    step_y = CNT_Y + Cm(0.3)

    for i, (col, badge, ttl, lines) in enumerate(steps):
        x = ML + i * (step_w + Cm(0.4))
        rect(sl, x, step_y, step_w, step_h, col)
        if i < len(steps) - 1:
            bx_arr, tf_arr = tbx(sl, x + step_w, step_y + Cm(4.0), Cm(0.4), Cm(1.5))
            para(tf_arr, "▶", size=9, color=C_MUTED, align=PP_ALIGN.CENTER, font=FONT_BODY)
        txt_c = C_WHITE if col in (C_COBALT, C_SAGE) else C_INK
        bx_b, tf_b = tbx(sl, x + Cm(0.2), step_y + Cm(0.2), step_w - Cm(0.3), Cm(0.6))
        para(tf_b, badge, size=9, bold=False, color=txt_c, font=FONT_UI)
        bx_t, tf_t = tbx(sl, x + Cm(0.2), step_y + Cm(0.85), step_w - Cm(0.3), Cm(1.1))
        para(tf_t, ttl, size=11, bold=False, color=txt_c, font=FONT_TITLE)
        div_col = RGBColor(0xFF,0xFF,0xFF) if col in (C_COBALT, C_SAGE) else RGBColor(0xCC,0xC8,0xBF)
        rect(sl, x + Cm(0.2), step_y + Cm(2.0), step_w - Cm(0.4), Cm(0.04), div_col)
        bx_c, tf_c = tbx(sl, x + Cm(0.2), step_y + Cm(2.2), step_w - Cm(0.3), Cm(6.8))
        for ln in lines:
            para(tf_c, f"· {ln}", size=9.5, color=txt_c, space_before=2, font=FONT_BODY)

    bx_d, tf_d = tbx(sl, ML, CNT_Y + Cm(10.1), CW, Cm(0.9))
    para(tf_d, "※ 중복 음식 자동 제거 (동일 food_name·serving_size·calories 튜플 dedup) + 그램합 오류 시 kcal 자동 재계산 보정",
         size=9.5, color=C_MUTED, italic=True, font=FONT_BODY)
    return sl


def slide_06_ai_detail(prs):
    sl = blank_slide(prs)
    std_frame(sl, "02.  AI 음식 분석 기술", "AI 추론 결과 구조 & 사용자 검증 흐름",
              subtitle="JSON 응답 스키마 · 양 조절 슬라이더 · DB 저장 구조", num=6)

    lw = CW * 0.48
    rect(sl, ML, CNT_Y, lw, Cm(9.8), C_INK)
    bx_l, tf_l = tbx(sl, ML + Cm(0.4), CNT_Y + Cm(0.3), lw - Cm(0.6), Cm(9.2))
    code_lines = [
        {'text': '// GPT-4o Vision 응답 JSON 스키마', 'size': 9,   'color': C_MUTED_S, 'italic': True, 'font': FONT_BODY},
        {'text': '{',                                  'size': 9.5, 'color': C_WHITE,   'font': FONT_BODY},
        {'text': '  "foods": [',                       'size': 9.5, 'color': C_WHITE,   'font': FONT_BODY},
        {'text': '    {',                              'size': 9.5, 'color': C_WHITE,   'font': FONT_BODY},
        {'text': '      "food_name": "비빔밥",',        'size': 9.5, 'color': C_SAGE,   'font': FONT_BODY},
        {'text': '      "calories": 560,',             'size': 9.5, 'color': C_OCHRE,  'font': FONT_BODY},
        {'text': '      "carbs": 95,',                 'size': 9.5, 'color': C_LAVEN,  'font': FONT_BODY},
        {'text': '      "protein": 18,',               'size': 9.5, 'color': C_LAVEN,  'font': FONT_BODY},
        {'text': '      "fat": 9,',                    'size': 9.5, 'color': C_LAVEN,  'font': FONT_BODY},
        {'text': '      "serving_size": "1인분",',     'size': 9.5, 'color': C_WHITE,  'font': FONT_BODY},
        {'text': '      "weight_g": 450,',             'size': 9.5, 'color': C_WHITE,  'font': FONT_BODY},
        {'text': '      "confidence": 0.87,',          'size': 9.5, 'color': C_MUTED_S,'font': FONT_BODY},
        {'text': '      "debug": {',                   'size': 9.5, 'color': C_MUTED_S,'font': FONT_BODY},
        {'text': '        "step1": "그릇 지름 약 22cm"','size': 9,  'color': C_MUTED_S,'font': FONT_BODY},
        {'text': '        "step2": "볶음밥류 밀도 적용"','size': 9, 'color': C_MUTED_S,'font': FONT_BODY},
        {'text': '      }',                            'size': 9.5, 'color': C_MUTED_S,'font': FONT_BODY},
        {'text': '    }',                              'size': 9.5, 'color': C_WHITE,  'font': FONT_BODY},
        {'text': '  ]',                                'size': 9.5, 'color': C_WHITE,  'font': FONT_BODY},
        {'text': '}',                                  'size': 9.5, 'color': C_WHITE,  'font': FONT_BODY},
    ]
    para_multi(tf_l, code_lines)

    rw = CW - lw - Cm(0.5)
    rx = ML + lw + Cm(0.5)
    steps_r = [
        (C_LAVEN,  "① 분석 결과 표시",     ["AI가 인식한 음식 목록·칼로리 제시", "신뢰도(confidence) 수치 함께 표시"]),
        (C_COBALT, "② 양 조절 슬라이더",   ["½ / 1 / 1½ / 2인분 선택", "칼로리·탄단지 실시간 재계산"]),
        (C_SAGE,   "③ 공유 그룹 선택",      ["개인 하루로그 기본 선택", "참여 중인 일반 그룹에도 공유 가능"]),
        (C_PEACH,  "④ DB 저장",            ["meal_records + meal_group_shares", "Cloudinary CDN 썸네일 자동 생성"]),
    ]
    sh_h = Cm(2.2)
    for i, (col, ttl, ls) in enumerate(steps_r):
        sy = CNT_Y + i * (sh_h + Cm(0.25))
        txt = C_WHITE if col in (C_COBALT, C_SAGE) else C_INK
        card(sl, rx, sy, rw, sh_h, col, ttl,
             [{'text': l, 'size': 9.5, 'color': txt, 'font': FONT_BODY} for l in ls],
             title_color=txt, t_size=11.5)

    bx_b, tf_b = tbx(sl, ML, CNT_Y + Cm(10.1), CW, Cm(0.9))
    para(tf_b, "텍스트 퀵로그(QuickLogModal) 지원 — 사진 없이 음식을 직접 검색해 기록, image_url=\"\" sentinel 처리",
         size=9.5, color=C_MUTED, italic=True, font=FONT_BODY)
    return sl


def slide_07_group(prs):
    sl = blank_slide(prs)
    std_frame(sl, "03.  주요 기능", "그룹 공유 & 실시간 비교",
              subtitle="친구·가족과 함께하는 소셜 식단 기록", num=7)

    card_w = (CW - Cm(0.5) * 2) / 3
    cards_top = [
        (C_COBALT, "그룹 생성 & 입장",
         ["6자리 코드로 즉시 참가", "2~12명 인원 제한 설정", "그룹명·설정 오너 전용 변경"]),
        (C_SAGE,   "그룹 피드",
         ["멤버 식사 사진·칼로리 피드", "아침/점심/저녁/간식 태그", "기록 없는 멤버도 '기록 없음' 표시"]),
        (C_PEACH,  "이모티콘 & 댓글",
         ["👍 😋 🔥 💪 😭 반응 버튼", "낙관적 업데이트로 즉각 반영", "인스타 라이브 스타일 이모지 플로팅"]),
    ]
    for i, (col, ttl, ls) in enumerate(cards_top):
        x = ML + i * (card_w + Cm(0.5))
        txt = C_WHITE if col in (C_COBALT, C_SAGE) else C_INK
        card(sl, x, CNT_Y, card_w, Cm(4.0), col, ttl,
             [{'text': l, 'size': 10, 'color': txt, 'font': FONT_BODY} for l in ls],
             title_color=txt, t_size=12)

    rect(sl, ML, CNT_Y + Cm(4.3), CW, Cm(0.06), C_HAIRLINE)

    bx_sub, tf_sub = tbx(sl, ML, CNT_Y + Cm(4.55), CW, Cm(0.7))
    para(tf_sub, "랭킹 & 비교 탭", size=13, bold=False, color=C_INK, font=FONT_TITLE)

    rank_items = [
        (C_COBALT, "그룹 내 랭킹",      "하루 칼로리 달성률 순위 실시간 집계 · 1위는 코발트 강조 카드 표시"),
        (C_OCHRE,  "달성률 바 그래프",  "0~79%(세이지) · 80~99%(황토) · 100%(코발트) · 110%+(코랄) 동적 색상"),
        (C_LAVEN,  "개인 하루로그 분리","is_personal=true 그룹은 비교에서 제외 · 모든 식사 자동 공유"),
    ]
    rh = Cm(1.6)
    for i, (col, ttl, desc) in enumerate(rank_items):
        ry = CNT_Y + Cm(5.4) + i * (rh + Cm(0.2))
        rect(sl, ML, ry, Cm(0.5), rh, col)
        bx_r, tf_r = tbx(sl, ML + Cm(0.7), ry + Cm(0.05), CW - Cm(0.7), rh)
        para(tf_r, ttl, size=11, bold=False, color=C_INK, font=FONT_TITLE)
        para(tf_r, desc, size=10, color=C_BODY, space_before=1, font=FONT_BODY)

    bx_n, tf_n = tbx(sl, ML, CNT_Y + Cm(10.3), CW, Cm(0.9))
    para(tf_n, "※ KST 표시 수정(UTC+9h 변환) · 리액션 낙관적 업데이트 · 달력 레이아웃 고정(42일 그리드) 완료",
         size=9, color=C_MUTED, italic=True, font=FONT_BODY)
    return sl


def slide_08_personal(prs):
    sl = blank_slide(prs)
    std_frame(sl, "03.  주요 기능", "개인 맞춤 건강 관리",
              subtitle="분석 탭 · 체중 트래커 · AI 코치 · 월간 통계", num=8)

    col_w = (CW - Cm(0.5)) / 2
    cards = [
        (C_COBALT, "분석 탭 — 일별 / 주간 / 월간",
         ["목표 칼로리 대비 달성률 진행바",
          "탄단지 원형 차트 (탄=황토·단=세이지·지=코랄)",
          "개인 목표 탄단지 대비 % 표시 (targets 프롭)",
          "단백질 기준 109g 통일 (식약처 2020 개정)"]),
        (C_SAGE,   "체중 기록 트래커",
         ["SVG 꺾은선 그래프 · 90일 추이 시각화",
          "오늘 체중 입력 · 같은 날 덮어쓰기",
          "POST /users/me/weight → UNIQUE(user_id, date)",
          "기록/삭제 시 Toast 알림 표시"]),
        (C_LAVEN,  "AI 코치",
         ["GPT-4o-mini 기반 맞춤 식단·운동 루틴",
          "PNG 내보내기 (Canvas API)",
          "운동명 → 구글 검색 링크 자동 연결",
          "결과 localStorage 저장 · 재방문 시 복원"]),
        (C_OCHRE,  "월간 통계 (MonthlyStats)",
         ["달력 히트맵 — 달성률 구간별 색상",
          "기록없음(회색)→세이지→황토→코발트→코랄",
          "월간 요약 수치 (총 칼로리·평균·달성일수)",
          "Top 5 자주 먹은 음식 표시"]),
    ]
    h_card = Cm(4.5)
    for i, (col, ttl, ls) in enumerate(cards):
        cx = ML + (i % 2) * (col_w + Cm(0.5))
        cy = CNT_Y + (i // 2) * (h_card + Cm(0.35))
        txt = C_WHITE if col in (C_COBALT, C_SAGE) else C_INK
        card(sl, cx, cy, col_w, h_card, col, ttl,
             [{'text': f"· {l}", 'size': 10, 'color': txt, 'font': FONT_BODY} for l in ls],
             title_color=txt, t_size=12)

    bx_b, tf_b = tbx(sl, ML, CNT_Y + Cm(9.65), CW, Cm(1.0))
    para(tf_b, "추가:  즐겨찾기 백엔드 DB 동기화 · 식사 알림(Notification API) · 식사 기록 수정(오너 전용) · 하루 식단 요약 카드 PNG 내보내기",
         size=9.5, color=C_MUTED, italic=True, font=FONT_BODY)
    return sl


def slide_09_demo(prs):
    """시연 전환 슬라이드"""
    sl = blank_slide(prs)
    bg(sl, C_COBALT)

    # 왼쪽 오프화이트 패널 (좁게)
    rect(sl, Cm(0), Cm(0), Cm(6), SH, C_CANVAS)

    # 좌측 — 챕터 번호
    bx_n, tf_n = tbx(sl, Cm(0.5), Cm(7.0), Cm(5.0), Cm(5.0))
    para(tf_n, "04", size=72, bold=False, color=C_COBALT, font=FONT_TITLE)
    para(tf_n, "DEMO", size=16, bold=False, color=C_MUTED, space_before=2, font=FONT_UI)

    # 우측 — 메인 메시지
    bx_t, tf_t = tbx(sl, Cm(7.5), Cm(5.5), Cm(24), Cm(3.5))
    para(tf_t, "프로그램 시연", size=40, bold=False, color=C_WHITE, font=FONT_TITLE)

    rect(sl, Cm(7.5), Cm(9.5), Cm(22), Cm(0.08), C_PEACH)

    bx_s, tf_s = tbx(sl, Cm(7.5), Cm(10.0), Cm(24), Cm(2.0))
    para(tf_s, "직접 화면을 보며 주요 기능을 설명드리겠습니다",
         size=16, bold=False, color=C_LAVEN, font=FONT_UI)

    # 시연 항목
    demo_items = [
        "·  음식 사진 촬영 → AI 3단계 분석 → 저장",
        "·  그룹 피드 · 이모티콘 반응 · 랭킹 확인",
        "·  분석 탭 — 일별 영양소 · 체중 트래커 · AI 코치",
    ]
    bx_l, tf_l = tbx(sl, Cm(7.5), Cm(12.3), Cm(24), Cm(3.5))
    for item in demo_items:
        para(tf_l, item, size=13, color=C_WHITE, space_before=4, font=FONT_BODY)

    # 하단 페이지 번호
    bx_p, tf_p = tbx(sl, SW - Cm(4), SH - Cm(1.2), Cm(3.5), Cm(0.8))
    para(tf_p, f"9  /  {TOTAL_SLIDES}", size=9, color=C_WHITE,
         align=PP_ALIGN.RIGHT, font=FONT_BODY)
    return sl


def slide_10_arch(prs):
    sl = blank_slide(prs)
    std_frame(sl, "05.  기술 아키텍처", "시스템 아키텍처 & 기술 스택",
              subtitle="Next.js App Router · FastAPI · Supabase · Cloudinary · GPT-4o", num=10)

    layer_w = (CW - Cm(0.5) * 2) / 3
    layers = [
        (C_LAVEN, "프론트엔드",
         [("Next.js 14",    "App Router · SSR/CSR 혼합"),
          ("TypeScript",    "엄격 타입 체계"),
          ("Tailwind CSS",  "먹로그 디자인 토큰"),
          ("Zustand",       "전역 상태 관리"),
          ("TanStack Query","서버 상태 캐시"),
          ("Axios",         "camelCase↔snake_case 자동 변환"),
          ("Canvas API",    "이미지 내보내기"),
          ("Lucide React",  "아이콘 시스템")]),
        (C_COBALT, "백엔드",
         [("FastAPI",      "async Python REST API"),
          ("SQLAlchemy",   "async ORM"),
          ("Pydantic v2",  "요청/응답 스키마 검증"),
          ("OAuth2.0",     "Google 소셜 로그인"),
          ("JWT",          "7일 만료 액세스 토큰"),
          ("Cloudinary",   "이미지 업로드 + CDN"),
          ("OpenAI SDK",   "GPT-4o Vision / GPT-4o-mini"),
          ("Uvicorn",      "ASGI 서버")]),
        (C_SAGE, "인프라 / DB",
         [("Supabase",       "PostgreSQL · Auth · Realtime"),
          ("Vercel",         "프론트엔드 자동 배포"),
          ("Railway",        "FastAPI 백엔드 배포"),
          ("meal_records",   "음식 기록 핵심 테이블"),
          ("weight_records", "체중 기록 (UNIQUE per day)"),
          ("user_favorite_foods", "즐겨찾기 영구 저장"),
          ("meal_group_shares",   "그룹 공유 매핑"),
          ("SQLite",         "개발 환경 (Supabase IPv6 우회)")]),
    ]

    lh = Cm(10.0)
    for i, (col, ttl, rows) in enumerate(layers):
        lx = ML + i * (layer_w + Cm(0.5))
        txt = C_WHITE if col in (C_COBALT, C_SAGE) else C_INK
        rect(sl, lx, CNT_Y, layer_w, lh, col)
        bx_t, tf_t = tbx(sl, lx + Cm(0.3), CNT_Y + Cm(0.15), layer_w - Cm(0.4), Cm(0.65))
        para(tf_t, ttl, size=12, bold=False, color=txt, font=FONT_TITLE)
        div_col = RGBColor(0xFF,0xFF,0xFF) if col in (C_COBALT, C_SAGE) else RGBColor(0xBB,0xBA,0xB0)
        rect(sl, lx + Cm(0.3), CNT_Y + Cm(0.88), layer_w - Cm(0.6), Cm(0.04), div_col)
        bx_c, tf_c = tbx(sl, lx + Cm(0.3), CNT_Y + Cm(1.0), layer_w - Cm(0.4), lh - Cm(1.2))
        for pkg, desc in rows:
            if tf_c.paragraphs[0].text != '':
                p = tf_c.add_paragraph()
            else:
                p = tf_c.paragraphs[0]
            r1 = p.add_run(); r1.text = f"{pkg:<18}"
            r1.font.name = FONT_BODY; r1.font.size = Pt(9.5)
            r1.font.bold = True; r1.font.color.rgb = txt
            r2 = p.add_run(); r2.text = desc
            r2.font.name = FONT_BODY; r2.font.size = Pt(9)
            r2.font.color.rgb = txt
            p.space_before = Pt(2.5)

    bx_b, tf_b = tbx(sl, ML, CNT_Y + Cm(10.3), CW, Cm(0.85))
    para(tf_b, "Axios 인터셉터: 모든 요청/응답에서 snake_case ↔ camelCase 자동 변환 · DevSidebar(Ctrl+F11)에서 실시간 네트워크·AI 추론 로그 확인 가능",
         size=9.5, color=C_MUTED, italic=True, font=FONT_BODY)
    return sl


def slide_11_progress(prs):
    sl = blank_slide(prs)
    std_frame(sl, "06.  개발 성과 & 향후 계획", "Phase 1 → Phase 2 완성 내역",
              subtitle="MVP 핵심 루프 완성 후 기능 강화 및 품질 개선", num=11)

    lw = CW * 0.47
    rw = CW - lw - Cm(0.5)
    rx = ML + lw + Cm(0.5)

    rect(sl, ML, CNT_Y, lw, Cm(10.0), C_COBALT)
    bx_l, tf_l = tbx(sl, ML + Cm(0.35), CNT_Y + Cm(0.2), lw - Cm(0.5), Cm(0.6))
    para(tf_l, "Phase 1  —  MVP  (기말 과제 제출 범위)", size=11.5, bold=False, color=C_WHITE, font=FONT_TITLE)
    rect(sl, ML + Cm(0.35), CNT_Y + Cm(0.88), lw - Cm(0.6), Cm(0.04), RGBColor(0xFF,0xFF,0xFF))
    bx_l2, tf_l2 = tbx(sl, ML + Cm(0.35), CNT_Y + Cm(1.0), lw - Cm(0.5), Cm(8.7))
    for ln in [
        "✓  Google OAuth 2.0 로그인",
        "✓  GPT-4o Vision AI 3단계 칼로리 분석",
        "✓  음식 정보 수정 + 양 조절 슬라이더",
        "✓  날짜별 개인 로그 (칼로리 진행바)",
        "✓  그룹 생성·참가 (2~12명 코드 방식)",
        "✓  그룹 피드 + 이모지 반응 + 댓글",
        "✓  그룹원 실시간 칼로리 비교·랭킹",
        "✓  Mifflin-St Jeor 개인 목표 칼로리",
        "✓  AI 보정 (그램합 오류 → kcal 재계산)",
        "✓  AI 중복 음식 자동 제거 (dedup)",
        "✓  그룹 설정 (이름·인원 오너 전용)",
        "✓  DevSidebar 개발자 모드",
        "✓  KST 시간 수정 · 달력 레이아웃 안정화",
    ]:
        para(tf_l2, ln, size=9.5, color=C_WHITE, space_before=2, font=FONT_BODY)

    rect(sl, rx, CNT_Y, rw, Cm(10.0), C_SAGE)
    bx_r, tf_r = tbx(sl, rx + Cm(0.35), CNT_Y + Cm(0.2), rw - Cm(0.5), Cm(0.6))
    para(tf_r, "Phase 2  —  기능 강화", size=11.5, bold=False, color=C_WHITE, font=FONT_TITLE)
    rect(sl, rx + Cm(0.35), CNT_Y + Cm(0.88), rw - Cm(0.6), Cm(0.04), RGBColor(0xFF,0xFF,0xFF))
    bx_r2, tf_r2 = tbx(sl, rx + Cm(0.35), CNT_Y + Cm(1.0), rw - Cm(0.5), Cm(8.7))
    for ln in [
        "✓  즐겨찾기 (로컬→백엔드 DB 동기화)",
        "✓  식사 알림 (Notification API)",
        "✓  식사 기록 수정 (오너 전용)",
        "✓  AI 코치 (식단·운동 + PNG 내보내기)",
        "✓  텍스트 퀵로그 (QuickLogModal)",
        "✓  체중 기록 트래커 (SVG 그래프)",
        "✓  월간 통계 (히트맵 + Top5)",
        "✓  Toast 알림 시스템",
        "✓  영양소 개인 목표 대비 % 표시",
        "✓  즐겨찾기 mapFood 버그 수정",
        "✓  퀵로그 image_url 500 에러 수정",
        "✓  AI 코치 결과 localStorage 저장",
    ]:
        para(tf_r2, ln, size=9.5, color=C_WHITE, space_before=2, font=FONT_BODY)

    bx_b, tf_b = tbx(sl, ML, CNT_Y + Cm(10.3), CW, Cm(0.85))
    para(tf_b, "총 완성: Phase 1 — 15개 · Phase 2 — 25개  |  미구현: RefreshToken · 그룹 채팅(WebSocket) · Rate Limiting · 그룹장 양도",
         size=9.5, color=C_MUTED, italic=True, font=FONT_BODY)
    return sl


def slide_12_future(prs):
    sl = blank_slide(prs)
    std_frame(sl, "06.  개발 성과 & 향후 계획", "스마트폰 앱 출시 및 Phase 3 로드맵",
              subtitle="현재 웹 앱 → 향후 React Native 기반 iOS / Android 앱으로 전환", num=12)

    lw = CW * 0.55
    rw = CW - lw - Cm(0.6)
    rx = ML + lw + Cm(0.6)

    bx_l, tf_l = tbx(sl, ML, CNT_Y, lw, Cm(2.5))
    para(tf_l, "스마트폰 앱 전환 계획", size=14, bold=False, color=C_COBALT, font=FONT_TITLE)
    para(tf_l, "현재 먹로그는 모바일 최적화 웹 앱으로 구현되어 있으며,\n추후 React Native(Expo) 기반 네이티브 앱으로 전환 후 App Store · Play Store에 정식 배포할 예정입니다.",
         size=10.5, color=C_BODY, space_before=4, font=FONT_BODY)

    phases = [
        (C_COBALT, "1단계 — 웹 앱 완성 (현재)",
         "Next.js 14 기반 PWA · 모바일 브라우저 최적화 · Vercel 배포 중"),
        (C_OCHRE,  "2단계 — React Native 전환",
         "Expo 프레임워크 도입 · 네이티브 카메라 API · 푸시 알림 · 오프라인 동작"),
        (C_SAGE,   "3단계 — 앱 스토어 출시",
         "iOS App Store · Google Play 정식 등록 · 온보딩 흐름 최적화"),
    ]
    ph_h = Cm(1.8)
    for i, (col, ttl, desc) in enumerate(phases):
        py = CNT_Y + Cm(2.8) + i * (ph_h + Cm(0.25))
        txt = C_WHITE if col in (C_COBALT, C_SAGE) else C_INK
        card(sl, ML, py, lw, ph_h, col, ttl,
             [{'text': desc, 'size': 9.5, 'color': txt, 'font': FONT_BODY}],
             title_color=txt, t_size=11)

    bx_rd, tf_rd = tbx(sl, ML, CNT_Y + Cm(8.5), lw, Cm(2.5))
    para(tf_rd, "Phase 3 기술 로드맵", size=12, bold=False, color=C_INK, font=FONT_TITLE)
    for ln in [
        "· YOLO 기반 다중 음식 객체 동시 탐지",
        "· 인바디 사진 OCR 체성분 자동 분석",
        "· 그룹 챌린지 시스템 (7일 달성 뱃지)",
        "· Supabase Realtime 이모지 실시간 반영",
        "· RefreshToken 도입 (현재 7일 만료 후 재로그인)",
    ]:
        para(tf_rd, ln, size=10, color=C_BODY, space_before=2, font=FONT_BODY)

    rect(sl, rx, CNT_Y, rw, Cm(11.0), C_SURF)
    bx_logo_t, tf_logo_t = tbx(sl, rx + Cm(0.3), CNT_Y + Cm(0.3), rw - Cm(0.5), Cm(0.8))
    para(tf_logo_t, "앱 아이콘 (확정)", size=11, bold=False, color=C_COBALT, font=FONT_TITLE)

    if os.path.exists(LOGO_PATH):
        logo_size = rw - Cm(1.6)
        logo_x = rx + (rw - logo_size) / 2
        sl.shapes.add_picture(LOGO_PATH, logo_x, CNT_Y + Cm(1.2), logo_size, logo_size)

    bx_logo_d, tf_logo_d = tbx(sl, rx + Cm(0.3), CNT_Y + Cm(8.4), rw - Cm(0.5), Cm(2.0))
    para(tf_logo_d, "귀여운 동물 캐릭터 아이콘으로\n친근하고 가벼운 앱 이미지 강조",
         size=9.5, color=C_MUTED, align=PP_ALIGN.CENTER, font=FONT_BODY)
    para(tf_logo_d, "iOS / Android 동시 출시 예정",
         size=10, bold=False, color=C_COBALT, space_before=4,
         align=PP_ALIGN.CENTER, font=FONT_TITLE)
    return sl


def slide_13_closing(prs):
    sl = blank_slide(prs)
    bg(sl, C_COBALT)
    rect(sl, Cm(0), Cm(0), SW * 0.58, SH, C_CANVAS)

    bx_m, tf_m = tbx(sl, Cm(2.2), Cm(3.5), Cm(17), Cm(4.0))
    para(tf_m, "감사합니다", size=48, bold=False, color=C_COBALT, font=FONT_TITLE)

    rect(sl, Cm(2.2), Cm(7.7), Cm(14), Cm(0.09), C_PEACH)

    bx_s, tf_s = tbx(sl, Cm(2.2), Cm(8.2), Cm(17), Cm(3.0))
    para(tf_s, "먹로그 — 사진 한 장으로 완성하는 AI 식단 기록",
         size=14, bold=False, color=C_INK, font=FONT_TITLE)
    para(tf_s, "AI와 머신러닝 기말고사 과제  ·  2026. 06. 22",
         size=11, color=C_MUTED, space_before=6, font=FONT_BODY)
    para(tf_s, "팀 먹로그  ·  25학번 이건우  ·  23학번 윤서영",
         size=11, color=C_MUTED, space_before=3, font=FONT_BODY)

    bx_k, tf_k = tbx(sl, Cm(2.2), Cm(11.5), Cm(17), Cm(4.5))
    para(tf_k, "핵심 성과 요약", size=12, bold=False, color=C_COBALT, font=FONT_TITLE)
    for k in [
        "Phase 1 + Phase 2 완성 — 40개 기능 구현",
        "GPT-4o Vision 3단계 물리 기반 칼로리 추론",
        "그룹 소셜 피드 + 실시간 비교 + 이모지 반응",
        "체중 트래커 · AI 코치 · 월간 히트맵 통계",
        "추후 iOS / Android 네이티브 앱 출시 예정",
    ]:
        para(tf_k, f"  ·  {k}", size=10.5, color=C_BODY, space_before=3, font=FONT_BODY)

    if os.path.exists(LOGO_PATH):
        sl.shapes.add_picture(LOGO_PATH, Cm(22.0), Cm(3.5), Cm(9.0), Cm(9.0))

    bx_n, tf_n = tbx(sl, SW - Cm(4), SH - Cm(1.0), Cm(3.5), Cm(0.8))
    para(tf_n, f"{TOTAL_SLIDES}  /  {TOTAL_SLIDES}", size=9, bold=False,
         color=C_WHITE, align=PP_ALIGN.RIGHT, font=FONT_BODY)
    return sl


# ═══════════════════════════════════════════════════════════════════════
# 메인
# ═══════════════════════════════════════════════════════════════════════

def main():
    prs = prs_new()
    funcs = [
        ("01/13  타이틀",             slide_01_title),
        ("02/13  목차",               slide_02_toc),
        ("03/13  문제 정의",          slide_03_problem),
        ("04/13  핵심 아이디어",      slide_04_solution),
        ("05/13  AI 3단계 추론 개요", slide_05_ai_overview),
        ("06/13  AI 분석 상세",       slide_06_ai_detail),
        ("07/13  그룹 기능",          slide_07_group),
        ("08/13  개인 건강 관리",     slide_08_personal),
        ("09/13  [시연 전환]",        slide_09_demo),
        ("10/13  기술 아키텍처",      slide_10_arch),
        ("11/13  개발 성과",          slide_11_progress),
        ("12/13  향후 계획",          slide_12_future),
        ("13/13  마무리",             slide_13_closing),
    ]
    for label, fn in funcs:
        fn(prs)
        print(f"  {label}")

    out = r"c:\Desktop\문서\AI_GUNWOO\Muklog\먹로그_발표자료_v1.6.0.pptx"
    prs.save(out)
    print("saved:", out)

if __name__ == "__main__":
    main()
